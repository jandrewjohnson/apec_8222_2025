from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    # Initialize Presentation
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Integrating Machine Learning into Causal Inference"
    subtitle.text = "Beyond Prediction: Tools for High-Dimensional Econometrics\nBased on: Arif, S. (2025)"

    # --- Slide 2: The Econometric Problem ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "The Econometric Problem: Functional Form & Dimensionality"
    
    tf = content.text_frame
    tf.text = "The Limitation of OLS"
    p = tf.add_paragraph()
    p.text = "Standard identification requires conditioning on confounders (Z)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "In high-dimensional datasets, relationships between Z and Y are complex/non-linear."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Risk of specification search ('p-hacking') or omitted variable bias."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The ML Opportunity"
    p = tf.add_paragraph()
    p.text = "Standard ML minimizes MSE for prediction but is biased for parameters."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Causal ML learns 'nuisance functions' flexibly to isolate structural parameters."
    p.level = 1

    slide.notes_slide.notes_text_frame.text = (
        "As economists, we are comfortable with identification strategies. "
        "However, we often assume linear separability in our controls. "
        "Causal ML allows us to relax the assumption of functional form for control variables."
    )

    # --- Slide 3: Identification remains "King" ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Identification remains 'King'"
    
    tf = content.text_frame
    tf.text = "Algorithms ≠ Identification"
    p = tf.add_paragraph()
    p.text = "No algorithm corrects for fundamental endogeneity without a valid strategy (IV, DAGs)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The Workflow:"
    p = tf.add_paragraph()
    p.text = "1. Structural Model: Define causal graph (DAG) or Potential Outcomes."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Identification: Establish backdoor criterion or valid instrument."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Estimation: Use ML to estimate conditional expectations E[Y|Z] non-parametrically."
    p.level = 1

    slide.notes_slide.notes_text_frame.text = (
        "ML is not a magic wand for bad research design. If you have unobserved confounding, "
        "a Random Forest will just overfit the bias. We still need rigorous identification."
    )

    # --- Slide 4: Double Machine Learning (DML) ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Double Machine Learning (DML)"
    
    tf = content.text_frame
    tf.text = "The 'Frisch-Waugh-Lovell' of ML"
    p = tf.add_paragraph()
    p.text = "Theoretical Basis: Chernozhukov et al. (2018)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Mechanism:"
    p = tf.add_paragraph()
    p.text = "1. Partialling Out: Predict Y from Z and D from Z."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Residuals: Calculate residuals for Y and D."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Inference: Regress residuals of Y on residuals of D."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Why it works:"
    p = tf.add_paragraph()
    p.text = "Cross-fitting: Splits sample to avoid overfitting bias."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Achieves root-N consistency for treatment coefficient."
    p.level = 1

    slide.notes_slide.notes_text_frame.text = (
        "Application: Estimating demand elasticity with high-dimensional consumer characteristics."
    )

    # --- Slide 5: TMLE ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Targeted Maximum Likelihood Estimation (TMLE)"
    
    tf = content.text_frame
    tf.text = "Efficiency in Binary Treatment Evaluation"
    p = tf.add_paragraph()
    p.text = "Concept: A doubly robust estimator (Van der Laan & Rubin, 2006)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Different from DML: Updates initial outcome estimate using a 'clever covariate' derived from propensity score."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Properties:"
    p = tf.add_paragraph()
    p.text = "Double Robustness: Consistent if either propensity score OR outcome model is correct."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Efficiency: Asymptotically efficient (lowest variance)."
    p.level = 1

    slide.notes_slide.notes_text_frame.text = (
        "Application: Program evaluation (e.g., job training) where selection bias is driven by complex observables."
    )

    # --- Slide 6: Deep IV ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Deep Instrumental Variables (Deep IV)"
    
    tf = content.text_frame
    tf.text = "Handling Non-Linearity in 2SLS"
    p = tf.add_paragraph()
    p.text = "Problem: If 1st stage relationship (D on Z) is non-linear, standard 2SLS is weak."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "The Solution:"
    p = tf.add_paragraph()
    p.text = "Use Deep Neural Networks to model instrument compliance function g(Z)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Maintains exclusion restriction (Z ⊥ ε)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Deep IV (Hartford et al., 2017):"
    p = tf.add_paragraph()
    p.text = "Projects inputs into a latent distribution to recover structural function."
    p.level = 1

    slide.notes_slide.notes_text_frame.text = (
        "Application: Demand estimation where instrument affects price in a highly non-linear way."
    )

    # --- Slide 7: Causal Forests ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Causal Forests & Heterogeneity"
    
    tf = content.text_frame
    tf.text = "Uncovering 'For Whom' the Policy Works"
    p = tf.add_paragraph()
    p.text = "Theoretical Basis: Athey & Imbens (2016); Wager & Athey (2018)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Goal:"
    p = tf.add_paragraph()
    p.text = "Estimate Conditional Average Treatment Effects (CATE): τ(x)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Mechanism:"
    p = tf.add_paragraph()
    p.text = "Adapts Random Forests to maximize heterogeneity in effects."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "'Honest' Splitting: Separate subsamples for tree structure vs. estimation."
    p.level = 1

    slide.notes_slide.notes_text_frame.text = (
        "Application: Targeted marketing or personalized medicine—identifying subgroups."
    )

    # --- Slide 8: Case Study ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Case Study: Non-Linearities (Ecological Analogy)"
    
    tf = content.text_frame
    tf.text = "Scenario: Estimating effect of environmental input (Depth) on output."
    
    p = tf.add_paragraph()
    p.text = "Standard OLS: Assumes constant marginal effect."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Causal Forest Result:"
    p = tf.add_paragraph()
    p.text = "Identified regime shifts: Effect positive in Region A, zero in Region B."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Driven by interacting stressors."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Implication:"
    p = tf.add_paragraph()
    p.text = "Standard ATE would mask these offsetting effects, leading to suboptimal allocation."
    p.level = 1

    slide.notes_slide.notes_text_frame.text = (
        "In an economic context, this is akin to a minimum wage study finding positive effects "
        "in concentrated labor markets but negative effects in competitive ones."
    )

    # --- Slide 9: Comparison Table ---
    slide_layout = prs.slide_layouts[5] # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Choosing the Right Estimator"
    
    # Add Table
    rows = 5
    cols = 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9.0)
    height = Inches(0.8)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(4.0)
    
    # Headers
    table.cell(0, 0).text = "Econometric Challenge"
    table.cell(0, 1).text = "Recommended Method"
    table.cell(0, 2).text = "Key Advantage"
    
    # Row 1
    table.cell(1, 0).text = "High-Dim Controls (N < P)"
    table.cell(1, 1).text = "Double ML (DML)"
    table.cell(1, 2).text = "Handles complex nuisance parameters; familiar regression output."
    
    # Row 2
    table.cell(2, 0).text = "Binary Policy / RCT"
    table.cell(2, 1).text = "TMLE"
    table.cell(2, 2).text = "Maximum efficiency; Doubly Robust protection."
    
    # Row 3
    table.cell(3, 0).text = "Endogeneity + Nonlinear IV"
    table.cell(3, 1).text = "Deep IV"
    table.cell(3, 2).text = "Recovers structural parameters when 1st stage is complex."
    
    # Row 4
    table.cell(4, 0).text = "Heterogeneity / Targeting"
    table.cell(4, 1).text = "Causal Forests"
    table.cell(4, 2).text = "Data-driven discovery of subgroups (CATE)."

    # --- Slide 10: Conclusion ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Conclusion"
    
    tf = content.text_frame
    tf.text = "Integrating ML into the Economist’s Toolkit"
    
    p = tf.add_paragraph()
    p.text = "1. Robustness: Reduces reliance on arbitrary parametric assumptions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Validation: Methods backed by rigorous asymptotic theory."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. Actionable Insight: Moving from ATE to CATE allows for optimal policy targeting."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Final Thought:"
    p = tf.add_paragraph()
    p.text = "Causal ML automates the 'statistical clean-up', allowing economists to focus on identification."
    p.level = 1

    # Save
    prs.save('Causal_ML_for_Economists.pptx')

create_presentation()